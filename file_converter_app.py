import streamlit as st
import os
import io
import base64
import json
from PyPDF2 import PdfReader
from docx import Document
import pptx
import csv

def extract_text_from_file(file):
    text = ""
    file_type = file.name.split(".")[-1].lower()

    try:
        if file_type == "pdf":
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() or ""
        elif file_type == "docx":
            doc = Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_type == "pptx":
            prs = pptx.Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_type == "csv":
            decoded = file.getvalue().decode("utf-8").splitlines()
            reader = csv.reader(decoded)
            for row in reader:
                text += ", ".join(row) + "\n"
        elif file_type == "txt":
            text += file.getvalue().decode("utf-8")
        else:
            text = f"[Unsupported file type: {file_type}]"
    except Exception as e:
        text = f"[Error reading file {file.name}: {str(e)}]"

    return text.strip()

def main():
    st.set_page_config(page_title="File to Text Converter", layout="centered")
    st.title("📄 File to Text Converter")
    st.subheader("Convert your documents to a single file for AI training")

    uploaded_files = st.file_uploader(
        "Select files to convert",
        accept_multiple_files=True,
        type=["pdf", "docx", "pptx", "csv", "txt"]
    )

    export_format = st.radio("Choose export format", ["Plain Text (.txt)", "Structured JSONL (.jsonl)"])

    if st.button("Convert Files 🚀"):
        if uploaded_files:
            all_texts = []
            for file in uploaded_files:
                text = extract_text_from_file(file)
                all_texts.append({"filename": file.name, "content": text})

            if export_format == "Plain Text (.txt)":
                combined = "\n\n-----\n\n".join([f"--- {doc['filename']} ---\n{doc['content']}" for doc in all_texts])
                st.download_button(
                    label="Download .txt File",
                    data=combined,
                    file_name="all_content.txt",
                    mime="text/plain"
                )
                st.text_area("Preview (first 5000 chars)", combined[:5000])
            else:
                jsonl_data = "\n".join([json.dumps(doc) for doc in all_texts])
                st.download_button(
                    label="Download .jsonl File",
                    data=jsonl_data,
                    file_name="all_content.jsonl",
                    mime="application/json"
                )
                st.text_area("Preview (first 5000 chars)", jsonl_data[:5000])
        else:
            st.warning("Please upload files first!")

if __name__ == "__main__":
    main()